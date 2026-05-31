import os

def get_pdf_text(file_path: str) -> str:
    """Extracts all text from a document (PDF, DOCX, PPTX).

    Args:
        file_path: Absolute path to the file.
    """
    text = ""
    try:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".docx":
            import docx
            doc = docx.Document(file_path)
            for p in doc.paragraphs:
                if p.text:
                    text += p.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            text += cell.text + " "
                    text += "\n"
        elif ext == ".pptx":
            import pptx
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
        else:
            import PyPDF2
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
    except Exception as e:
        return f"Error extracting text: {str(e)}"
    
    return text.strip()


if __name__ == "__main__":
    # Test with a sample file if needed
    import sys
    if len(sys.argv) > 1:
        print(get_pdf_text(sys.argv[1]))
