import os
try:
    import importlib.metadata as metadata
except ImportError:
    import importlib_metadata as metadata

from typing import Optional
from google.cloud import translate_v3 as translate
from google.adk.tools.tool_context import ToolContext

async def adaptive_translate_tool(
    file_path: str, 
    target_language_code: str, 
    source_language_code: Optional[str] = "en-US",
    glossary_id: Optional[str] = None,
    customized_attribution: Optional[str] = "NO_ATTRIBUTION",
    tool_context: ToolContext = None
) -> dict:
    """Translates a document while preserving layout using Google Cloud Translation API Advanced.

    Args:
        file_path: Absolute path to the PDF document to translate.
        target_language_code: The BCP-47 language code to translate into (e.g., 'es', 'fr').
        source_language_code: The BCP-47 language code of the source document.
        glossary_id: Optional ID of the glossary to use for consistent terminology.
        customized_attribution: Attribution watermark setting. Defaults to "NO_ATTRIBUTION" to remove watermark.
        tool_context: The ADK ToolContext for accessing state and project info.
    """
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    location = os.getenv("TRANSLATION_LOCATION", "us-central1")
    
    client = translate.TranslationServiceClient()
    parent = f"projects/{project_id}/locations/{location}"

    # Load the document content
    try:
        with open(file_path, "rb") as document_file:
            content = document_file.read()
    except Exception as e:
        return {"status": "error", "message": f"Failed to read file: {str(e)}"}

    # Determine extension and MIME type
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    elif ext == ".pptx":
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        mime_type = "application/pdf"
        ext = ".pdf"

    if ext == ".pdf":
        # Check if document has searchable text using PyMuPDF
        import fitz
        has_text = False
        try:
            doc = fitz.open(file_path)
            for page in doc:
                if len(page.get_text("text").strip()) > 50:
                    has_text = True
                    break
            doc.close()
        except Exception as e:
            print(f"Failed to check text content with PyMuPDF: {e}")
            # Assume it has text if we can't open it with fitz (let Translation API try)
            has_text = True
        if not has_text:
            print("Document has very little or no searchable text. Triggering full Gemini translation fallback...")
            try:
                from PIL import Image
                import io
                from google import genai
                from google.genai import types
                
                doc = fitz.open(file_path)
                out_doc = fitz.open() # New empty PDF
                
                project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
                client_gen = genai.Client(vertexai=True, project=project_id, location="global")
                
                for page_num in range(len(doc)):
                    print(f"Fallback processing page {page_num} with Gemini...")
                    page = doc[page_num]
                    pix = page.get_pixmap(dpi=300)
                    img_bytes = pix.tobytes("png")
                    
                    generator_prompt = f"""
                    Translate ALL text on this page into {target_language_code}.
                    Recreate the page visual exactly, preserving the layout, styles, colors, and any charts or diagrams.
                    The output must be a high-quality image containing the translated content.
                    Ensure all text is translated, leaving nothing in the original language.
                    """
                    
                    gen_response = client_gen.models.generate_content(
                        model="publishers/google/models/gemini-3.1-flash-image",
                        contents=[
                            types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                            types.Part.from_text(text=generator_prompt)
                        ]
                    )
                    
                    new_img_bytes = None
                    for part in gen_response.candidates[0].content.parts:
                        try:
                            if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                                new_img_bytes = part.inline_data.data
                                break
                        except AttributeError:
                            pass
                            
                    if new_img_bytes:
                        img = Image.open(io.BytesIO(new_img_bytes))
                        img_pdf_bytes = io.BytesIO()
                        img.save(img_pdf_bytes, format="PDF")
                        img_doc = fitz.open("pdf", img_pdf_bytes.getvalue())
                        out_doc.insert_pdf(img_doc)
                    else:
                        print(f"Failed to generate translated page {page_num} with Gemini.")
                        # Insert original page as fallback
                        out_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)
                        
                output_file_path = file_path.replace(".pdf", f"_{target_language_code}.pdf")
                out_doc.save(output_file_path)
                out_doc.close()
                doc.close()
                
                return {
                    "status": "success", 
                    "output_file": output_file_path,
                    "detected_language": source_language_code
                }
            except Exception as fallback_err:
                print(f"Full Gemini fallback failed: {fallback_err}")
                return {"status": "error", "message": f"Translation failed on zero-text document: {str(fallback_err)}"}

    # Proceed with standard Translation API if text is present or if docx/pptx
    document_input_config = {
        "content": content,
        "mime_type": mime_type,
    }

    # Prepare the request
    request = {
        "parent": parent,
        "target_language_code": target_language_code,
        "source_language_code": source_language_code,
        "document_input_config": document_input_config,
    }

    if customized_attribution:
        request["customized_attribution"] = customized_attribution

    if glossary_id:
        glossary_config = {
            "glossary": f"projects/{project_id}/locations/{location}/glossaries/{glossary_id}"
        }
        request["glossary_config"] = glossary_config

    try:
        response = client.translate_document(request=request)
        translated_content = response.document_translation.byte_stream_outputs[0]
        
        # Save output file
        output_file_path = file_path.replace(ext, f"_{target_language_code}{ext}")
        with open(output_file_path, "wb") as f:
            f.write(translated_content)
        
        if ext == ".pdf":
            # Post-process images to translate text within them
            try:
                print(f"Processing images in {output_file_path}...")
                await localize_images_in_pdf(file_path, output_file_path, target_language_code)
            except Exception as e:
                print(f"Image translation failed: {e}")
                
            # Refine table layouts and typography (Solution B & C)
            try:
                print(f"Refining table layout and typography in {output_file_path}...")
                await refine_table_layout_and_typography(file_path, output_file_path, target_language_code)
            except Exception as e:
                print(f"Table refinement failed: {e}")
                
            # Embed native CJK font subset if target is Japanese, Chinese, or Korean
            if target_language_code in ["ja", "zh", "ko"]:
                try:
                    print(f"Embedding native CJK font subset into {output_file_path}...")
                    doc_cjk = fitz.open(output_file_path)
                    cjk_paths = [
                        "/System/Library/Fonts/STHeiti Light.ttc",
                        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
                        "/usr/share/fonts/truetype/droid/DroidSansFallback.ttf",
                        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                        "C:\\Windows\\Fonts\\msgothic.ttc"
                    ]
                    cjk_font_path = None
                    for path in cjk_paths:
                        if os.path.exists(path):
                            cjk_font_path = path
                            break
                            
                    if cjk_font_path:
                        for page in doc_cjk:
                            page.insert_font(fontname="cjk", fontfile=cjk_font_path)
                        temp_cjk_path = output_file_path + ".cjk.pdf"
                        doc_cjk.save(temp_cjk_path, incremental=False)
                        doc_cjk.close()
                        import shutil
                        shutil.move(temp_cjk_path, output_file_path)
                        print("Native CJK font subset embedded successfully!")
                    else:
                        doc_cjk.close()
                        print("No CJK system font found to embed.")
                except Exception as cjk_err:
                    print(f"CJK font embedding failed: {cjk_err}")
        elif ext == ".docx":
            try:
                print(f"Processing images in DOCX {output_file_path}...")
                await localize_images_in_docx(file_path, output_file_path, target_language_code)
            except Exception as e:
                print(f"DOCX Image translation failed: {e}")
        elif ext == ".pptx":
            try:
                print(f"Processing images in PPTX {output_file_path}...")
                await localize_images_in_pptx(file_path, output_file_path, target_language_code)
            except Exception as e:
                print(f"PPTX Image translation failed: {e}")
            
        return {
            "status": "success", 
            "output_file": output_file_path,
            "detected_language": response.document_translation.detected_language_code
        }
    except Exception as e:
        return {"status": "error", "message": f"Translation failed: {str(e)}"}

async def localize_images_in_docx(original_path: str, translated_path: str, target_lang: str):
    import docx
    from google import genai
    from google.genai import types
    import io
    import asyncio

    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    client = genai.Client(vertexai=True, project=project_id, location="global")

    doc = docx.Document(translated_path)
    
    async def call_gemini_async(client_obj, bytes_data, mime_type, prompt_text):
        try:
            def sync_call():
                return client_obj.models.generate_content(
                    model="publishers/google/models/gemini-3.1-flash-image",
                    contents=[
                        types.Part.from_bytes(data=bytes_data, mime_type=mime_type),
                        types.Part.from_text(text=prompt_text)
                    ]
                )
            return await asyncio.to_thread(sync_call)
        except Exception as e:
            print(f"Async Gemini call via thread failed: {e}")
            return None

    tasks = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_part.content_type:
            img_part = rel.target_part
            image_bytes = img_part.blob
            mime = img_part.content_type

            generator_prompt = f"""
            Translate ALL text within this image into {target_lang}.
            It is CRITICAL that every single word, label, title, and legend item is translated to {target_lang}.
            Do NOT leave any text in English.
            Generate a new image that is identical in style, layout, colors, and data presentation as the input image, but with the fully translated text.
            Ensure high visual fidelity and crisp text.
            """
            tasks.append((img_part, image_bytes, mime, generator_prompt))

    if not tasks:
        print("No images found in DOCX to localize.")
        return

    print(f"Executing {len(tasks)} DOCX image translation tasks in parallel...")
    coroutines = [call_gemini_async(client, t[1], t[2], t[3]) for t in tasks]
    results = await asyncio.gather(*coroutines)

    for i, result in enumerate(results):
        if not result:
            continue
        new_img_bytes = None
        for part in result.candidates[0].content.parts:
            try:
                if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                    new_img_bytes = part.inline_data.data
                    break
            except AttributeError:
                pass
        if new_img_bytes:
            img_part = tasks[i][0]
            img_part._blob = new_img_bytes

    doc.save(translated_path)
    print(f"Successfully localized images in DOCX at {translated_path}")

async def localize_images_in_pptx(original_path: str, translated_path: str, target_lang: str):
    import pptx
    from google import genai
    from google.genai import types
    import io
    import asyncio

    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    client = genai.Client(vertexai=True, project=project_id, location="global")

    prs = pptx.Presentation(translated_path)
    
    async def call_gemini_async(client_obj, bytes_data, mime_type, prompt_text):
        try:
            def sync_call():
                return client_obj.models.generate_content(
                    model="publishers/google/models/gemini-3.1-flash-image",
                    contents=[
                        types.Part.from_bytes(data=bytes_data, mime_type=mime_type),
                        types.Part.from_text(text=prompt_text)
                    ]
                )
            return await asyncio.to_thread(sync_call)
        except Exception as e:
            print(f"Async Gemini call via thread failed: {e}")
            return None

    tasks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13: # mso_shape_type PICTURE = 13
                image = shape.image
                image_bytes = image.blob
                mime = f"image/{image.ext}" if image.ext else "image/png"
                generator_prompt = f"""
                Translate ALL text within this image into {target_lang}.
                It is CRITICAL that every single word, label, title, and legend item is translated to {target_lang}.
                Do NOT leave any text in English.
                Generate a new image that is identical in style, layout, colors, and data presentation as the input image, but with the fully translated text.
                Ensure high visual fidelity and crisp text.
                """
                tasks.append((image, image_bytes, mime, generator_prompt))

    if not tasks:
        print("No images found in PPTX to localize.")
        return

    print(f"Executing {len(tasks)} PPTX image translation tasks in parallel...")
    coroutines = [call_gemini_async(client, t[1], t[2], t[3]) for t in tasks]
    results = await asyncio.gather(*coroutines)

    for i, result in enumerate(results):
        if not result:
            continue
        new_img_bytes = None
        for part in result.candidates[0].content.parts:
            try:
                if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                    new_img_bytes = part.inline_data.data
                    break
            except AttributeError:
                pass
        if new_img_bytes:
            image_obj = tasks[i][0]
            image_obj._blob = new_img_bytes

    prs.save(translated_path)
    print(f"Successfully localized images in PPTX at {translated_path}")


async def localize_images_in_pdf(original_pdf_path: str, translated_pdf_path: str, target_lang: str):
    """Finds charts/images in the PDF, uses Gemini to regenerate them with translated text,
    and inserts them back into the PDF. Uses parallel processing for speed.
    """
    import fitz  # PyMuPDF
    from google import genai
    from google.genai import types
    import json
    import io
    import asyncio
    from PIL import Image
    
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    client = genai.Client(vertexai=True, project=project_id, location="global")
    
    doc_orig = fitz.open(original_pdf_path)
    doc_trans = fitz.open(translated_pdf_path)
    
    tasks = []
    
    # Helper to call Gemini asynchronously using threads
    async def call_gemini_async(client_obj, bytes_data, mime_type, prompt_text):
        try:
            def sync_call():
                return client_obj.models.generate_content(
                    model="publishers/google/models/gemini-3.1-flash-image",
                    contents=[
                        types.Part.from_bytes(data=bytes_data, mime_type=mime_type),
                        types.Part.from_text(text=prompt_text)
                    ]
                )
            # Run the synchronous call in a separate thread to avoid blocking
            response = await asyncio.to_thread(sync_call)
            return response
        except Exception as e:
            print(f"Async Gemini call via thread failed: {e}")
            return None

    # Step 1: Gather all tasks across all pages
    for page_num in range(len(doc_orig)):
        page_orig = doc_orig[page_num]
        images = page_orig.get_images(full=True)
        image_info = page_orig.get_image_info(hashes=True)
        
        print(f"Scanning page {page_num} for tasks...")
        page_tasks_added = 0
        used_bboxes = set()  # Track already matched bboxes on this page to resolve collisions natively
        
        for img in images:
            xref = img[0]
            width = img[2]
            height = img[3]
            
            if width < 100 or height < 100:
                continue
                
            matching_info = None
            # Original width/height matching, but tracking already allocated bboxes to resolve collisions in relative order
            for info in image_info:
                bbox_tuple = tuple(info['bbox'])
                if bbox_tuple in used_bboxes:
                    continue
                if info['width'] == width and info['height'] == height:
                    matching_info = info
                    used_bboxes.add(bbox_tuple)
                    break
                    
            if not matching_info:
                continue
                
            bbox = matching_info['bbox']
            
            try:
                base_image = doc_orig.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]
            except Exception as e:
                print(f"Failed to extract image {xref}: {e}")
                continue
                
            generator_prompt = f"""
            Analyze this image for any English text (such as labels, titles, legends, or words).
            If the image contains NO English text at all (or only contains numbers, symbols, or decorative graphics), you MUST respond with exactly the text: NO_TEXT
            Otherwise, if it contains English text, translate all text to {target_lang} and generate a new image that is identical in style, layout, colors, and data presentation as the input image, but with the fully translated text.
            The output image MUST be generated to match or scale nicely to {width}x{height} pixels.
            Ensure high visual fidelity and crisp text.
            """
            
            tasks.append({
                "type": "image",
                "page_num": page_num,
                "xref": xref,
                "width": width,
                "height": height,
                "bbox": bbox,
                "bytes": image_bytes,
                "mime": f"image/{ext}",
                "prompt": generator_prompt
            })
            page_tasks_added += 1
            
        # Fallback if no large images found on this page AND no tables detected
        try:
            tables = page_orig.find_tables()
            has_tables = len(tables.tables) > 0
        except:
            has_tables = False
            
        if page_tasks_added == 0 and not has_tables:
            print(f"No large images or tables on page {page_num}. Checking for vector drawings...")
            try:
                drawings = page_orig.get_drawings()
                page_width = page_orig.rect.width
                page_height = page_orig.rect.height
                
                chart_drawings = []
                for d in drawings:
                    d_rect = fitz.Rect(d["rect"])
                    # Ignore decorative page-sized background blocks
                    if d_rect.width > 0.9 * page_width and d_rect.height > 0.9 * page_height:
                        continue
                    chart_drawings.append(d_rect)
                
                if chart_drawings:
                    print(f"Found {len(chart_drawings)} vector drawings. Computing intelligent targeted bounding box...")
                    # Start with drawings boundary
                    chart_rect = chart_drawings[0]
                    for r in chart_drawings[1:]:
                        chart_rect |= r
                    
                    # Dynamically find text blocks on this page that are likely chart labels/legends
                    try:
                        blocks = page_orig.get_text("blocks")
                        max_dist = 50.0  # maximum points away from drawings core
                        
                        for b in blocks:
                            b_rect = fitz.Rect(b[:4])
                            b_text = b[4].strip()
                            
                            # Filter out dense body paragraphs or headers based on character length and height
                            if len(b_text) > 120 or b_rect.height > 40:
                                continue
                                
                            # Calculate horizontal distance
                            if b_rect.x1 < chart_rect.x0:
                                dx = chart_rect.x0 - b_rect.x1
                            elif b_rect.x0 > chart_rect.x1:
                                dx = b_rect.x0 - chart_rect.x1
                            else:
                                dx = 0.0
                                
                            # Calculate vertical distance
                            if b_rect.y1 < chart_rect.y0:
                                dy = chart_rect.y0 - b_rect.y1
                            elif b_rect.y0 > chart_rect.y1:
                                dy = b_rect.y0 - chart_rect.y1
                            else:
                                dy = 0.0
                                
                            dist = max(dx, dy)
                            if dist <= max_dist:
                                chart_rect |= b_rect
                                print(f"Dynamically included chart label block: '{b_text[:30]}...' with bbox {b_rect}")
                    except Exception as text_err:
                        print(f"Failed to dynamically cluster text blocks: {text_err}")
                    
                    # Apply uniform 15-point safety padding to prevent edge clipping
                    x0 = max(0.0, chart_rect.x0 - 15.0)
                    y0 = max(0.0, chart_rect.y0 - 15.0)
                    x1 = min(page_width, chart_rect.x1 + 15.0)
                    y1 = min(page_height, chart_rect.y1 + 15.0)
                    
                    rect = fitz.Rect(x0, y0, x1, y1)
                    print(f"Targeted fallback bounding box for page {page_num}: {rect}")
                    
                    pix = page_orig.get_pixmap(clip=rect, dpi=300)
                    img_bytes = pix.tobytes("png")
                    
                    generator_prompt = f"""
                    Analyze this chart image for any English text (such as labels, titles, legends, or words).
                    If the image contains NO English text at all (or only contains numbers, symbols, or decorative graphics), you MUST respond with exactly the text: NO_TEXT
                    Otherwise, if it contains English text, translate all text within this chart image into {target_lang}.
                    Generate a new image that is identical in style, layout, colors, and data presentation as the input image, but with the fully translated text.
                    Ensure high visual fidelity and crisp text.
                    """
                    
                    tasks.append({
                        "type": "fallback",
                        "page_num": page_num,
                        "rect": rect,
                        "bytes": img_bytes,
                        "mime": "image/png",
                        "prompt": generator_prompt
                    })
                    print(f"Added targeted fallback task for page {page_num} with rect {rect}")
                else:
                    print(f"No charts/drawings detected on page {page_num}. Keeping page text-only.")
            except Exception as e:
                print(f"Failed to create targeted fallback task for page {page_num}: {e}")

    # Step 2: Execute all tasks in parallel
    async def run_tasks():
        coroutines = []
        for t in tasks:
            coroutines.append(call_gemini_async(client, t["bytes"], t["mime"], t["prompt"]))
        return await asyncio.gather(*coroutines)
        
    print(f"Executing {len(tasks)} translation tasks in parallel...")
    results = await run_tasks()
    print("All parallel tasks completed.")

    # Step 3: Apply results sequentially to the document
    for i, result in enumerate(results):
        task = tasks[i]
        if not result:
            print(f"Skipping failed result for task {i}")
            continue
            
        # Check if Gemini requested to skip replacement due to lack of text
        is_no_text = False
        try:
            for part in result.candidates[0].content.parts:
                if part.text and "NO_TEXT" in part.text:
                    is_no_text = True
                    break
        except:
            pass
            
        if is_no_text:
            print(f"Task {i} on page {task['page_num']} has no translatable English text. Skipping replacement to preserve original graphics.")
            continue
            
        new_img_bytes = None
        for part in result.candidates[0].content.parts:
            try:
                if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                    new_img_bytes = part.inline_data.data
                    break
            except AttributeError:
                pass
                
        if not new_img_bytes:
            print(f"No image data in result for task {i}")
            continue
            
        page_trans = doc_trans[task["page_num"]]
        
        try:
            # Resize if needed
            if task["type"] == "image":
                width, height = task["width"], task["height"]
                rect = fitz.Rect(task["bbox"])
            else:
                rect = task["rect"]
                width, height = int(rect.width), int(rect.height) # Approximation or read from image
                
            # Force dimensions
            from PIL import Image
            img = Image.open(io.BytesIO(new_img_bytes))
            img_resized = img.resize((width, height), Image.Resampling.LANCZOS)
            buffer = io.BytesIO()
            img_resized.save(buffer, format="PNG")
            final_bytes = buffer.getvalue()
            
            # Dynamic text protection for image type
            if task["type"] == "image":
                try:
                    words = page_trans.get_text("words")
                    for w in words:
                        w_rect = fitz.Rect(w[:4])
                        if w_rect.intersects(rect) and w_rect.x0 > rect.x0 + rect.width * 0.7:
                            rect.x1 = min(rect.x1, w_rect.x0 - 5)
                except:
                    pass
                    
            page_trans.insert_image(rect, stream=final_bytes, keep_proportion=True, overlay=True)
            print(f"Applied result for task {i} on page {task['page_num']}")
            
        except Exception as e:
            print(f"Failed to apply result for task {i}: {e}")

    temp_path = translated_pdf_path + ".tmp"
    doc_trans.save(temp_path, incremental=False, encryption=fitz.PDF_ENCRYPT_KEEP)
    doc_trans.close()
    doc_orig.close()
    
    import shutil
    shutil.move(temp_path, translated_pdf_path)
    print(f"Successfully saved clean PDF via temp file to {translated_pdf_path}")

async def fix_tables_in_pdf(original_pdf_path: str, translated_pdf_path: str, target_lang: str):
    """Uses Gemini to calculate corrected coordinates for table text to prevent column bleed,
    by detecting table bounding boxes on the ORIGINAL PDF and applying fixes to the TRANSLATED PDF.
    """
    import fitz
    from google import genai
    from google.genai import types
    import json
    import re
    import asyncio
    
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    client = genai.Client(vertexai=True, project=project_id, location="global")
    
    doc_orig = fitz.open(original_pdf_path)
    doc_trans = fitz.open(translated_pdf_path)
    
    for page_num in range(len(doc_orig)):
        page_orig = doc_orig[page_num]
        page_trans = doc_trans[page_num]
        
        # Detect tables on the ORIGINAL page where structure is clean
        try:
            tables = page_orig.find_tables()
            table_bboxes = []
            
            if len(tables.tables) == 0:
                print(f"PyMuPDF found no tables on page {page_num}. Trying AI detection...")
                try:
                    # Render whole page to find table area
                    pix = page_orig.get_pixmap(dpi=300)
                    img_bytes = pix.tobytes("png")
                    
                    detect_prompt = """
                    Identify the bounding box of the main table on this page. 
                    Return the coordinates as a JSON object with normalized values from 0 to 1000:
                    {"ymin": 200, "xmin": 50, "ymax": 800, "xmax": 950}
                    If no table exists, return {"no_table": true}.
                    """
                    
                    def sync_call_detect():
                        return client.models.generate_content(
                            model="publishers/google/models/gemini-3.5-flash",
                            contents=[
                                types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                                types.Part.from_text(text=detect_prompt)
                            ]
                        )
                    response = await asyncio.to_thread(sync_call_detect)
                    
                    json_match = re.search(r'\{.*\}', response.text, re.DOTALL)
                    if json_match:
                        data = json.loads(json_match.group(0))
                        if "no_table" not in data:
                            width = page_orig.rect.width
                            height = page_orig.rect.height
                            rect = fitz.Rect(
                                (data["xmin"] / 1000.0) * width,
                                (data["ymin"] / 1000.0) * height,
                                (data["xmax"] / 1000.0) * width,
                                (data["ymax"] / 1000.0) * height
                            )
                            table_bboxes.append(rect)
                            print(f"AI detected table on page {page_num} with bbox {rect}")
                except Exception as detect_err:
                    print(f"AI table detection failed: {detect_err}")
            else:
                for table in tables.tables:
                    table_bboxes.append(fitz.Rect(table.bbox))
                    
            if not table_bboxes:
                continue
                
            for tab_idx, rect in enumerate(table_bboxes):
                print(f"Fixing table {tab_idx} on page {page_num} with bbox {rect}")
                
                # Render the messed up table region from the TRANSLATED page
                pix = page_trans.get_pixmap(clip=rect, dpi=300)
                img_bytes = pix.tobytes("png")
                
                prompt = f"""
                In this table image, the text columns are overlapping due to word length expansion.
                Your job is to extract all text strings and calculate the corrected X and Y coordinates for each string so they sit properly centered in their respective columns without overlapping.
                
                Return the output as a JSON array of objects ONLY:
                [
                    {{"text": "Frances Gullefant", "x": 150, "y": 200}},
                    ...
                ]
                Coordinates MUST be normalized from 0 to 1000 relative to this image area (0 is left/top, 1000 is right/bottom).
                Ensure data in the same row has exactly the same Y coordinate.
                """
                
                def sync_call():
                    return client.models.generate_content(
                        model="publishers/google/models/gemini-3.5-flash",
                        contents=[
                            types.Part.from_bytes(data=img_bytes, mime_type="image/png"),
                            types.Part.from_text(text=prompt)
                        ]
                    )
                
                response = await asyncio.to_thread(sync_call)
                
                # Parse JSON
                json_match = re.search(r'\[.*\]', response.text, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group(0))
                    
                    # 1. Wipe out old table area with white box on TRANSLATED page
                    page_trans.draw_rect(rect, color=(1, 1, 1), fill=(1, 1, 1))
                    
                    # 2. Insert new text at clean coordinates
                    for item in data:
                        text = item["text"]
                        x_norm = item["x"]
                        y_norm = item["y"]
                        
                        x_pdf = rect.x0 + (x_norm / 1000.0) * rect.width
                        y_pdf = rect.y0 + (y_norm / 1000.0) * rect.height
                        
                        page_trans.insert_text(fitz.Point(x_pdf, y_pdf), text, fontsize=9, color=(0, 0, 0))
                        
                    print(f"Successfully redrew table {tab_idx} on page {page_num}")
                else:
                    print(f"No valid JSON array found in response for table {tab_idx}")
        except Exception as e:
            print(f"Failed to fix table on page {page_num}: {e}")
            continue
            
    doc_trans.save(translated_pdf_path, incremental=True, encryption=fitz.PDF_ENCRYPT_KEEP)
    doc_trans.close()
    doc_orig.close()

def inspect_translation_proof(file_path: str, glossary_terms: list, source_content: Optional[str] = None) -> dict:
    """Uses Gemini to inspect the translated PDF and verify glossary enforcement with source context.
    
    Args:
        file_path: Absolute path to the translated PDF document.
        glossary_terms: List of terms expected to be preserved/translated correctly.
        source_content: Optional text of the source document to prevent false-negatives.
    """
    if not glossary_terms:
        return {
            "status": "success", 
            "audit_report": "No glossary terms were provided for this translation. Glossary enforcement check skipped."
        }

    from google import genai
    from google.genai import types
    
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    location = "global"
    
    client = genai.Client(vertexai=True, project=project_id, location=location)
    
    try:
        with open(file_path, "rb") as f:
            pdf_bytes = f.read()
            
        prompt = f"""You are an AI Auditor for Financial Translations. 
Review the attached translated PDF document. 

### CONTEXT
- **Expected Glossary Terms**: {", ".join(glossary_terms)}
- **Original Source Text (for reference)**: {source_content[:5000] if source_content else "Not provided"}

### YOUR MISSION
1. Verify if the 'Expected Glossary Terms' were strictly followed in the translated PDF.
2. **CRITICAL**: If a glossary term is not present in the 'Original Source Text', IGNORE it. Do not report it as a failure if it wasn't in the source.
3. For terms that ARE in the source, confirm if they were correctly preserved or localized in the output.
"""
        response = client.models.generate_content(
            model="gemini-2.5-pro",
            contents=[
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_bytes(data=pdf_bytes, mime_type="application/pdf"),
                        types.Part.from_text(text=prompt)
                    ]
                )
            ]
        )
        return {"status": "success", "audit_report": response.text}
    except Exception as e:
        return {"status": "error", "message": f"AI Audit failed: {str(e)}"}

def read_csv_glossary(file_path: str) -> dict:
    """Reads a custom CSV glossary to understand terminology enforcement.
    
    Args:
        file_path: Absolute path to the CSV glossary file.
    """
    import csv
    if not os.path.exists(file_path):
        return {"status": "error", "message": f"Glossary file not found at {file_path}"}
    
    terms = []
    try:
        with open(file_path, mode='r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for row in reader:
                if row:
                    terms.append(" -> ".join(row))
        return {"status": "success", "terms": terms}
    except Exception as e:
        return {"status": "error", "message": str(e)}

def save_dynamic_glossary(terms: list, filename: str = "dynamic_glossary.csv") -> dict:
    """Saves a list of terminology mappings to a temporary CSV glossary.
    
    Args:
        terms: List of strings in 'SourceTerm,TargetTerm' format.
        filename: Name of the temporary glossary file.
    """
    import csv
    base_dir = os.path.dirname(os.path.abspath(__file__))
    glossary_dir = os.path.join(base_dir, "glossaries")
    os.makedirs(glossary_dir, exist_ok=True)
    
    file_path = os.path.join(glossary_dir, filename)
    try:
        with open(file_path, mode='w', encoding='utf-8', newline='') as f:
            writer = csv.writer(f)
            for term in terms:
                if "," in term:
                    writer.writerow(term.split(",", 1))
                else:
                    writer.writerow([term, term]) # Preserve if no mapping
        return {"status": "success", "file_path": os.path.abspath(file_path)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

async def refine_table_layout_and_typography(original_pdf_path: str, translated_pdf_path: str, target_lang: str):
    """Refines table layouts and typography inside translated PDFs using targeted Gemini Image-to-Image Translation.
    Only runs on specific page regions where tables are detected. Uses parallel processing with a concurrency throttle.
    1. Identifies the precise table bounding box on the original page.
    2. Renders only that table bounding box as a high-res image.
    3. Translates all tables in parallel using gemini-3.1-flash-image (throttled to max 10 concurrent requests).
    4. Wipes out only the table bounding box on the translated page with a white mask.
    5. Overlays the fully translated table image exactly on top of the table region.
    """
    import fitz
    import os
    import io
    import asyncio
    from google import genai
    from google.genai import types
    
    project_id = os.getenv("GOOGLE_CLOUD_PROJECT", "uppdemos")
    client = genai.Client(vertexai=True, project=project_id, location="global")
    
    doc_orig = fitz.open(original_pdf_path)
    doc_trans = fitz.open(translated_pdf_path)
    
    # Define a semaphore to throttle concurrency (max 10 parallel requests to respect API quotas/sockets)
    sem = asyncio.Semaphore(10)
    
    # Helper to call Gemini asynchronously using threads and a semaphore throttle
    async def call_gemini_async(client_obj, bytes_data, prompt_text, semaphore):
        async with semaphore:
            try:
                def sync_call():
                    return client_obj.models.generate_content(
                        model="gemini-3.1-flash-image",
                        contents=[
                            types.Part.from_bytes(data=bytes_data, mime_type="image/png"),
                            types.Part.from_text(text=prompt_text)
                        ]
                    )
                # Run the synchronous call in a separate thread to avoid blocking
                response = await asyncio.to_thread(sync_call)
                return response
            except Exception as e:
                print(f"Async Gemini table call via thread failed: {e}")
                return None

    tasks = []
    
    # Step 1: Gather all table tasks across all pages
    for page_num in range(len(doc_orig)):
        page_orig = doc_orig[page_num]
        
        # Check if tables are present on this page
        try:
            tables = page_orig.find_tables()
            if len(tables.tables) == 0:
                continue
        except Exception as e:
            print(f"Table detection failed on page {page_num}: {e}")
            continue
            
        print(f"Page {page_num}: {len(tables.tables)} Table(s) detected for parallel refinement.")
        
        for tab_idx, table in enumerate(tables.tables):
            bbox = fitz.Rect(table.bbox)
            
            # Render high-resolution image of the targeted table region only
            pix = page_orig.get_pixmap(clip=bbox, dpi=300)
            img_bytes = pix.tobytes("png")
            
            prompt = f"""
            Translate ALL English text within this table image into the target language: '{target_lang}'.
            It is CRITICAL that every single word, label, title, header, and number is translated accurately.
            Generate a new image that is identical in style, layout, colors, grid lines, and structure as the input image, but with the fully translated text.
            Ensure high visual fidelity, crisp legible text, and correct alignment inside cells.
            """
            
            tasks.append({
                "page_num": page_num,
                "tab_idx": tab_idx,
                "bbox": bbox,
                "bytes": img_bytes,
                "prompt": prompt
            })

    if not tasks:
        print("No tables detected in document. Table refinement skipped.")
        doc_trans.close()
        doc_orig.close()
        return

    # Step 2: Execute all table translation tasks in parallel (throttled by the semaphore)
    async def run_tasks():
        coroutines = []
        for t in tasks:
            coroutines.append(call_gemini_async(client, t["bytes"], t["prompt"], sem))
        return await asyncio.gather(*coroutines)
        
    print(f"Executing {len(tasks)} table translation tasks in parallel (throttled)...")
    results = await run_tasks()
    print("All parallel table tasks completed.")
    
    modified = False
    
    # Step 3: Apply results sequentially to the document
    for i, response in enumerate(results):
        task = tasks[i]
        if not response:
            print(f"Skipping failed table result for table {task['tab_idx']} on page {task['page_num']}")
            continue
            
        # Iterate backwards through the response parts to grab the final/best generated image
        translated_img_bytes = None
        try:
            if response.candidates and response.candidates[0].content.parts:
                for part in reversed(response.candidates[0].content.parts):
                    if getattr(part, 'inline_data', None) and part.inline_data and part.inline_data.mime_type.startswith("image/"):
                        translated_img_bytes = part.inline_data.data
                        break
        except Exception as part_err:
            print(f"Error parsing response parts for table {task['tab_idx']} on page {task['page_num']}: {part_err}")
            
        if translated_img_bytes:
            print(f"Successfully received translated image for table {task['tab_idx']} on page {task['page_num']}!")
            page_trans = doc_trans[task["page_num"]]
            bbox = task["bbox"]
            
            # 1. Wipe out only the table region on the translated page
            page_trans.draw_rect(bbox, color=(1, 1, 1), fill=(1, 1, 1), overlay=True)
            
            # 2. Overlay the translated table image perfectly on top
            page_trans.insert_image(bbox, stream=translated_img_bytes, keep_proportion=False, overlay=True)
            
            modified = True
        else:
            print(f"Model failed to return a translated image for table {task['tab_idx']} on page {task['page_num']}.")

    if modified:
        temp_path = translated_pdf_path + ".tmp"
        doc_trans.save(temp_path, incremental=False, encryption=fitz.PDF_ENCRYPT_KEEP)
        doc_trans.close()
        doc_orig.close()
        
        import shutil
        shutil.move(temp_path, translated_pdf_path)
        print(f"Successfully saved refined table PDF to {translated_pdf_path}")
    else:
        doc_trans.close()
        doc_orig.close()

